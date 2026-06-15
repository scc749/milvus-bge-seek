"""Document parsing service for ingestion sources."""

from __future__ import annotations

import csv
import json
from pathlib import Path
from typing import Any


class DocumentParserService:
    """Parse local documents into normalized raw text records."""

    SIMPLE_TEXT_EXTENSIONS = {"md", "markdown", "txt", "text"}
    STRUCTURED_TEXT_EXTENSIONS = {"csv", "json", "html", "htm"}
    OFFICE_EXTENSIONS = {"pdf", "doc", "docx", "ppt", "pptx", "xls", "xlsx"}
    IMAGE_EXTENSIONS = {"jpg", "jpeg", "png", "gif", "bmp", "tiff", "webp"}
    AUDIO_EXTENSIONS = {"mp3", "wav", "m4a", "flac", "ogg"}
    VIDEO_EXTENSIONS = {"mp4", "mov", "avi", "mkv", "webm", "wmv", "flv"}

    def list_supported_extensions(self) -> dict[str, list[str]]:
        """Return supported extension groups for documentation and UI use."""

        return {
            "simple_text": sorted(self.SIMPLE_TEXT_EXTENSIONS),
            "structured_text": sorted(self.STRUCTURED_TEXT_EXTENSIONS),
            "office": sorted(self.OFFICE_EXTENSIONS),
            "image": sorted(self.IMAGE_EXTENSIONS),
            "audio": sorted(self.AUDIO_EXTENSIONS),
            "video": sorted(self.VIDEO_EXTENSIONS),
        }

    def parse_directory(self, directory_path: Path) -> list[dict[str, Any]]:
        """Parse every supported file in a directory tree."""

        documents: list[dict[str, Any]] = []
        for file_path in sorted(path for path in directory_path.rglob("*") if path.is_file()):
            parsed = self.parse_file(file_path, source_root=directory_path)
            if parsed is not None:
                documents.append(parsed)
        return documents

    def parse_file(
        self,
        file_path: Path,
        source_root: Path | None = None,
    ) -> dict[str, Any] | None:
        """Parse one local file into an ingestion-ready raw document."""

        extension = self._extension(file_path)
        parser_name = self._resolve_parser_name(extension)
        metadata = {
            "source_uri": str(file_path),
            "file_name": file_path.name,
            "file_extension": extension,
            "parser_name": parser_name,
        }
        if source_root is not None:
            metadata["source_root"] = str(source_root)

        if extension in self.SIMPLE_TEXT_EXTENSIONS:
            content = file_path.read_text(encoding="utf-8", errors="ignore")
        elif extension == "csv":
            content = self._csv_to_markdown(file_path)
        elif extension == "json":
            content = self._json_to_markdown(file_path)
        elif extension in {"html", "htm"}:
            content = self._html_to_text(file_path)
        elif extension == "pdf":
            content = self._parse_pdf(file_path)
        elif extension in {"doc", "docx"}:
            content = self._parse_word(file_path)
        elif extension in {"ppt", "pptx"}:
            content = self._parse_powerpoint(file_path)
        elif extension in {"xls", "xlsx"}:
            content = self._parse_spreadsheet(file_path)
        elif extension in self.IMAGE_EXTENSIONS:
            content = self._binary_placeholder("image", file_path)
        elif extension in self.AUDIO_EXTENSIONS:
            content = self._binary_placeholder("audio", file_path)
        elif extension in self.VIDEO_EXTENSIONS:
            content = self._binary_placeholder("video", file_path)
        else:
            return None

        return {
            "source_uri": str(file_path),
            "source_type": "file",
            "title": file_path.name,
            "content": content,
            "metadata": metadata,
        }

    def _resolve_parser_name(self, extension: str) -> str:
        if extension in self.SIMPLE_TEXT_EXTENSIONS | self.STRUCTURED_TEXT_EXTENSIONS:
            return "simple"
        if extension in self.OFFICE_EXTENSIONS:
            return "builtin"
        if extension in self.IMAGE_EXTENSIONS | self.AUDIO_EXTENSIONS | self.VIDEO_EXTENSIONS:
            return "multimodal_placeholder"
        return "unsupported"

    def _extension(self, file_path: Path) -> str:
        return file_path.suffix.lower().lstrip(".")

    def _csv_to_markdown(self, file_path: Path) -> str:
        with file_path.open("r", encoding="utf-8", errors="ignore", newline="") as handle:
            reader = csv.reader(handle)
            rows = list(reader)
        if not rows:
            return ""
        header = rows[0]
        lines = [
            "| " + " | ".join(header) + " |",
            "|" + "".join(" --- |" for _ in header),
        ]
        for row in rows[1:]:
            padded = [row[index] if index < len(row) else "" for index in range(len(header))]
            lines.append("| " + " | ".join(padded) + " |")
        return "\n".join(lines)

    def _json_to_markdown(self, file_path: Path) -> str:
        payload = json.loads(file_path.read_text(encoding="utf-8", errors="ignore"))
        return "```json\n" + json.dumps(payload, ensure_ascii=False, indent=2) + "\n```"

    def _html_to_text(self, file_path: Path) -> str:
        html_content = file_path.read_text(encoding="utf-8", errors="ignore")
        try:
            from bs4 import BeautifulSoup

            soup = BeautifulSoup(html_content, "html.parser")
            text = soup.get_text(separator="\n")
            return "\n".join(line.strip() for line in text.splitlines() if line.strip())
        except Exception:
            return html_content

    def _parse_pdf(self, file_path: Path) -> str:
        try:
            from pypdf import PdfReader

            reader = PdfReader(str(file_path))
            parts = [page.extract_text() or "" for page in reader.pages]
            return "\n\n".join(part for part in parts if part.strip())
        except Exception:
            return self._binary_placeholder("pdf", file_path)

    def _parse_word(self, file_path: Path) -> str:
        try:
            import docx2txt

            text = docx2txt.process(str(file_path))
            if text:
                return text
        except Exception:
            pass
        return self._try_unstructured(file_path, fallback_kind="word")

    def _parse_powerpoint(self, file_path: Path) -> str:
        if self._extension(file_path) == "pptx":
            try:
                from pptx import Presentation

                presentation = Presentation(str(file_path))
                slide_blocks: list[str] = []
                for slide_index, slide in enumerate(presentation.slides, start=1):
                    texts = [
                        shape.text.strip()
                        for shape in slide.shapes
                        if hasattr(shape, "text") and shape.text.strip()
                    ]
                    if texts:
                        slide_blocks.append(
                            f"## Slide {slide_index}\n" + "\n".join(f"- {text}" for text in texts)
                        )
                if slide_blocks:
                    return "\n\n".join(slide_blocks)
            except Exception:
                pass
        return self._try_unstructured(file_path, fallback_kind="powerpoint")

    def _parse_spreadsheet(self, file_path: Path) -> str:
        extension = self._extension(file_path)
        if extension == "xlsx":
            try:
                from openpyxl import load_workbook

                workbook = load_workbook(filename=str(file_path), data_only=True)
                sheet_blocks: list[str] = []
                for sheet in workbook.worksheets:
                    rows = list(sheet.iter_rows(values_only=True))
                    if not rows:
                        continue
                    header = [self._safe_cell(value) for value in rows[0]]
                    lines = [
                        f"## Sheet: {sheet.title}",
                        "| " + " | ".join(header) + " |",
                        "|" + "".join(" --- |" for _ in header),
                    ]
                    for row in rows[1:]:
                        lines.append(
                            "| "
                            + " | ".join(self._safe_cell(value) for value in row[: len(header)])
                            + " |"
                        )
                    sheet_blocks.append("\n".join(lines))
                if sheet_blocks:
                    return "\n\n".join(sheet_blocks)
            except Exception:
                pass
        return self._try_unstructured(file_path, fallback_kind="spreadsheet")

    def _try_unstructured(self, file_path: Path, fallback_kind: str) -> str:
        try:
            from langchain_community.document_loaders import UnstructuredFileLoader

            docs = UnstructuredFileLoader(str(file_path)).load()
            text = "\n\n".join(doc.page_content for doc in docs if doc.page_content.strip())
            if text:
                return text
        except Exception:
            pass
        return self._binary_placeholder(fallback_kind, file_path)

    def _binary_placeholder(self, kind: str, file_path: Path) -> str:
        return (
            f"[{kind.upper()} file]\n"
            f"name: {file_path.name}\n"
            f"path: {file_path}\n"
            "The current pipeline keeps this file in the ingestion flow and preserves metadata, "
            "but does not yet extract rich multimodal content."
        )

    def _safe_cell(self, value: Any) -> str:
        if value is None:
            return ""
        return str(value)
